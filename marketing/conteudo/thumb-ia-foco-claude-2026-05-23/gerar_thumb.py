from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import io, os

def px(n): return Emu(int(n * 9525))

# Caminhos
BASE = os.path.dirname(os.path.abspath(__file__))
FOTO_ORIG = os.path.join(BASE, "../../../identidade/IMG_8049.JPG")
FOTO_ROTACIONADA = os.path.join(BASE, "foto_rotacionada.jpg")
SAIDA = os.path.join(BASE, "thumb-ednei-ia.pptx")

# Cores
VERDE     = RGBColor(0x06, 0x44, 0x20)
DOURADO   = RGBColor(0xD4, 0xAF, 0x37)
BRANCO    = RGBColor(0xF8, 0xF5, 0xF0)
PRETO     = RGBColor(0x00, 0x00, 0x00)

W = 1280
H = 720

# ── 1. Rotacionar a foto ─────────────────────────────────────────────────────
img = Image.open(FOTO_ORIG)
img_rot = img.rotate(90, expand=True)   # 90° CCW → person upright
# Crop pra focar no rosto/torso (top-center do resultado)
iw, ih = img_rot.size
crop_w = int(iw * 0.55)
crop_h = ih
left = (iw - crop_w) // 2
img_crop = img_rot.crop((left, 0, left + crop_w, crop_h))
img_crop = img_crop.resize((560, 720), Image.LANCZOS)
img_crop.save(FOTO_ROTACIONADA, quality=92)
print(f"Foto processada: {img_crop.size}")

# ── 2. Criar apresentação ────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = px(W)
prs.slide_height = px(H)

slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

def add_rect(slide, x, y, w, h, fill_rgb=None, fill_alpha_hex=None, line=False):
    shape = slide.shapes.add_shape(1, px(x), px(y), px(w), px(h))
    shape.line.fill.background() if not line else None
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, font_size, bold, color, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Inter"
    return txBox

# ── Fundo verde ──────────────────────────────────────────────────────────────
bg = add_rect(slide, 0, 0, W, H, fill_rgb=VERDE)

# ── Foto à esquerda ───────────────────────────────────────────────────────────
pic = slide.shapes.add_picture(FOTO_ROTACIONADA, px(0), px(0), px(540), px(720))

# ── Overlay gradiente sobre a foto (retângulo com degradê simulado por sobreposição) ──
# Canva vai deixar editar isso como shape separada
ov = add_rect(slide, 300, 0, 280, 720)
ov.fill.solid()
ov.fill.fore_color.rgb = VERDE
ov.fill.fore_color.theme_color  # só pra garantir
# transparência parcial via XML
from lxml import etree
spPr = ov.fill._xPr
# Adiciona gradFill pra simular fade
grad_xml = '''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
  <a:gsLst>
    <a:gs pos="0">
      <a:srgbClr val="064420">
        <a:alpha val="0"/>
      </a:srgbClr>
    </a:gs>
    <a:gs pos="100000">
      <a:srgbClr val="064420">
        <a:alpha val="100000"/>
      </a:srgbClr>
    </a:gs>
  </a:gsLst>
  <a:lin ang="0" scaled="0"/>
</a:gradFill>'''

# Remove fill atual e insere gradFill
fill_el = ov._element.spPr
for child in list(fill_el):
    if 'Fill' in child.tag or 'fill' in child.tag.lower():
        fill_el.remove(child)
fill_el.append(etree.fromstring(grad_xml))

# ── Badge "IA NA PRÁTICA" ─────────────────────────────────────────────────────
badge_bg = add_rect(slide, 590, 90, 190, 30, fill_rgb=RGBColor(0x1a, 0x38, 0x1a))
add_text(slide, "  •  IA NA PRÁTICA", 590, 91, 190, 28,
         font_size=9, bold=False, color=DOURADO, align=PP_ALIGN.LEFT)

# ── Título ────────────────────────────────────────────────────────────────────
add_text(slide, "APRENDA", 590, 140, 640, 100,
         font_size=82, bold=True, color=BRANCO)

add_text(slide, "UMA IA", 590, 220, 640, 100,
         font_size=82, bold=True, color=DOURADO)

add_text(slide, "SÓ", 590, 300, 640, 100,
         font_size=82, bold=True, color=BRANCO)

add_text(slide, "ANTES DE APRENDER\nTODAS AS OUTRAS", 590, 400, 640, 80,
         font_size=24, bold=True, color=RGBColor(0xC0, 0xBC, 0xB7))

# ── Divisor dourado ───────────────────────────────────────────────────────────
add_rect(slide, 590, 495, 55, 3, fill_rgb=DOURADO)

# ── Bloco Claude ─────────────────────────────────────────────────────────────
# Quadrado dourado com "C"
claude_box = add_rect(slide, 590, 515, 52, 52, fill_rgb=DOURADO)
add_text(slide, "C", 594, 516, 44, 50,
         font_size=26, bold=True, color=VERDE, align=PP_ALIGN.CENTER)

# Texto "comece pelo / Claude"
add_text(slide, "COMECE PELO", 655, 518, 300, 22,
         font_size=10, bold=False, color=RGBColor(0x88, 0x88, 0x88))
add_text(slide, "Claude", 655, 538, 300, 32,
         font_size=22, bold=True, color=BRANCO)

# ── Salvar ────────────────────────────────────────────────────────────────────
prs.save(SAIDA)
print(f"\nArquivo gerado: {SAIDA}")
